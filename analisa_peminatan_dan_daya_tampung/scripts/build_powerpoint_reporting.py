import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def set_shape_flat_color(shape, rgb_color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_color
    shape.line.color.rgb = rgb_color

def create_header(slide, title_text, category_text="UNIVERSITAS SYIAH KUALA | LAPORAN STRATEGIS PMB"):
    # Header container
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.1))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

    p_cat = tf.paragraphs[0]
    p_cat.text = category_text.upper()
    p_cat.font.size = Pt(10)
    p_cat.font.bold = True
    p_cat.font.color.rgb = RGBColor(230, 81, 0) # Orange accent

    p_title = tf.add_paragraph()
    p_title.text = title_text
    p_title.font.size = Pt(22)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(27, 59, 111) # Navy Primary

def add_bullet(tf, bold_prefix, text_content, font_size=13, space_after=10):
    p = tf.add_paragraph()
    p.font.size = Pt(font_size)
    p.space_after = Pt(space_after)
    
    r_bold = p.add_run()
    r_bold.text = bold_prefix + " "
    r_bold.font.bold = True
    r_bold.font.color.rgb = RGBColor(33, 33, 33)

    r_norm = p.add_run()
    r_norm.text = text_content
    r_norm.font.color.rgb = RGBColor(66, 66, 66)

def main():
    base_dir = "/Users/auliamuzhaffar/Documents/maganghub"
    prs_dir = os.path.join(base_dir, "tugas-5", "analisa_peminatan_dan_daya_tampung")
    chart_dir = os.path.join(prs_dir, "grafik")
    out_pptx = os.path.join(prs_dir, "presentasi_analisa_peminatan_dan_daya_tampung.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # completely blank

    # -------------------------------------------------------------
    # SLIDE 1: COVER SLIDE
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    # Background banner
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_shape_flat_color(bg, RGBColor(248, 249, 250))
    bg.line.fill.background()

    # Left decorative bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
    set_shape_flat_color(bar, RGBColor(27, 59, 111))

    # Main Title box
    tb = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "TUGAS 05 - MAGANG (KHUSUS)"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = RGBColor(230, 81, 0)
    p0.space_after = Pt(10)

    p1 = tf.add_paragraph()
    p1.text = "ANALISA PEMINATAN DAN DAYA TAMPUNG\nPROGRAM STUDI UNIVERSITAS SYIAH KUALA"
    p1.font.size = Pt(30)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(27, 59, 111)
    p1.space_after = Pt(15)

    p2 = tf.add_paragraph()
    p2.text = "Evaluasi Tren Multi-Tahun (2022–2026), Rasio Keketatan Seleksi, Efisiensi Kapasitas Kuota, dan Kausalitas Lapangan sebagai Bahan Pertimbangan Strategis Penerimaan Mahasiswa Baru Menuju Rasionalisasi Kuota PTN-BH"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(85, 85, 85)
    p2.space_after = Pt(25)

    p3 = tf.add_paragraph()
    p3.text = "Disusun oleh: Tim Magang Universitas Syiah Kuala | Sumber Data Resmi: PMB USK 2022–2026"
    p3.font.size = Pt(12)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(120, 120, 120)

    # -------------------------------------------------------------
    # SLIDE 2: RINGKASAN EKSEKUTIF (SEGITIGA EMAS PMB)
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    create_header(s2, "Ringkasan Eksekutif: Membedah Segitiga Emas PMB USK (2022–2026)")
    
    # 3 Column Cards
    col_w = Inches(3.6)
    card_y = Inches(1.8)
    card_h = Inches(5.0)

    # Card 1: Peminat vs Kuota
    c1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), card_y, col_w, card_h)
    set_shape_flat_color(c1, RGBColor(255, 255, 255))
    c1.line.color.rgb = RGBColor(220, 220, 220)
    tf1 = c1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = Inches(0.25)
    p = tf1.paragraphs[0]
    p.text = "1. PEMINAT vs KUOTA"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(12)
    add_bullet(tf1, "Lonjakan Peminat:", "Total peminat USK naik dari 48.769 (2022) menjadi 68.010 (2026) setelah sempat menyentuh rekor 70.945 peminat di 2025.")
    add_bullet(tf1, "Ekspansi Kuota:", "Daya tampung diperluas agresif pasca PTN-BH sebesar +32,7% (dari 7.863 menjadi 10.435 kursi).")
    add_bullet(tf1, "Disparitas Keketatan:", "Farmasi terketat (39:1), sementara Budidaya Perairan peminatnya lebih sedikit dari kursi (0,91:1).")

    # Card 2: Realisasi Daftar Ulang
    c2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), card_y, col_w, card_h)
    set_shape_flat_color(c2, RGBColor(255, 255, 255))
    c2.line.color.rgb = RGBColor(220, 220, 220)
    tf2 = c2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = tf2.margin_right = tf2.margin_top = Inches(0.25)
    p = tf2.paragraphs[0]
    p.text = "2. REALISASI & KURSI KOSONG"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(230, 81, 0)
    p.space_after = Pt(12)
    add_bullet(tf2, "Pendaftar Ulang Naik:", "Mahasiswa masuk riil tumbuh dari 6.197 (2022) ke 8.440 mhs (2026) dengan keterisian rata-rata 80,9%.")
    add_bullet(tf2, "Beban 2.000 Bangku Kosong:", "USK secara persisten menyisakan ~2.000 kursi kosong setiap tahunnya akibat over-ekspansi kuota.")
    add_bullet(tf2, "Pemisahan Segmen:", "S1 Utama relatif sehat (84,7%), namun Vokasi D3 terpuruk di 51,3% dan PSDKU Gayo Lues hanya terisi 18,6%.")

    # Card 3: Kausalitas Lapangan
    c3 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), card_y, col_w, card_h)
    set_shape_flat_color(c3, RGBColor(255, 255, 255))
    c3.line.color.rgb = RGBColor(220, 220, 220)
    tf3 = c3.text_frame
    tf3.word_wrap = True
    tf3.margin_left = tf3.margin_right = tf3.margin_top = Inches(0.25)
    p = tf3.paragraphs[0]
    p.text = "3. PENYEBAB & SOLUSI"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(46, 125, 50)
    p.space_after = Pt(12)
    add_bullet(tf3, "Mengapa Prodi A Tinggi?", "Kepastian kerja ASN/PPPK (PGSD) dan pasar ners luar negeri (Keperawatan) mendorong permintaan tinggi.")
    add_bullet(tf3, "Mengapa Prodi B Rendah?", "Persepsi minim karir di daerah (Fisika, Budidaya Perairan) dan kanibalisasi internal Bisnis Digital.")
    add_bullet(tf3, "Kebocoran Mandiri & Talenta:", "Biaya IPI puluhan juta di Jalur Mandiri dan tiket cadangan Jalur Talenta menjadi sumber mundurnya 2.300+ calon.")

    # -------------------------------------------------------------
    # SLIDE 3: TREN MAKRO 5 TAHUN USK (GRAFIK 1)
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    create_header(s3, "Evaluasi Tren Makro USK 2022–2026: Kapasitas Kuota vs Realisasi Peminat")
    # Insert Chart
    img1 = os.path.join(chart_dir, "01_tren_makro_peminat_dt_du_usk.png")
    if os.path.exists(img1):
        s3.shapes.add_picture(img1, Inches(0.8), Inches(1.6), width=Inches(7.8))
    
    # Text commentary right
    tb = s3.shapes.add_textbox(Inches(8.9), Inches(1.6), Inches(3.8), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "POKOK TEMUAN MAKRO:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(10)
    add_bullet(tf, "1. Pertumbuhan Volume:", "Total peminat melonjak +39,5% selama 5 tahun. USK semakin diminati pendaftar secara regional dan nasional.")
    add_bullet(tf, "2. Lonjakan Kuota PTN-BH:", "Daya tampung dinaikkan dari 7.863 ke 10.435 kursi (+2.572 bangku) sejak transisi PTN-BH tahun 2024.")
    add_bullet(tf, "3. Jebakan 2.000 Kursi Kosong:", "Meskipun mahasiswa masuk meningkat, USK selalu gagal mengisi 1.995 s.d. 2.481 kursi setiap tahun.")
    add_bullet(tf, "4. Masalah Distribusi:", "Bukan kekurangan peminat, melainkan alokasi kuota berlebihan di prodi yang peminatnya sedikit.")

    # -------------------------------------------------------------
    # SLIDE 4: PETA KEKETATAN SELEKSI (GRAFIK 4)
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    create_header(s4, "Peta Tingkat Keketatan Seleksi 2026: Prodi Paling Favorit vs Paling Sepi")
    img4 = os.path.join(chart_dir, "04_rasio_keketatan_peminatan_vs_daya_tampung.png")
    if os.path.exists(img4):
        s4.shapes.add_picture(img4, Inches(0.8), Inches(1.6), width=Inches(8.0))
    
    tb = s4.shapes.add_textbox(Inches(9.1), Inches(1.6), Inches(3.6), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "FAKTA KEKETATAN SELEKSI:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(10)
    add_bullet(tf, "Top Favorit (Persaingan Sengit):", "Farmasi memimpin (39,2:1), Akuntansi Perpajakan D4 (22,4:1), Informatika (18,1:1), Psikologi (17,3:1), dan T. Perminyakan (15,5:1).")
    add_bullet(tf, "Anomali Kritis (Peminat < Kuota):", "Budidaya Perairan hanya memiliki rasio 0,91:1 (peminat 145 orang vs kuota 160 kursi). Secara matematis kuota mustahil terisi.")
    add_bullet(tf, "Prodi Rawan Sepi (<1,5:1):", "PSP Perikanan (1,07:1), Fisika (1,18:1), Pendidikan Fisika (1,29:1), dan Proteksi Tanaman (1,31:1).")

    # -------------------------------------------------------------
    # SLIDE 5: MATRIKS 4 KUADRAN STRATEGIS PORTOFOLIO PRODI (GRAFIK 10)
    # -------------------------------------------------------------
    s_kuadran = prs.slides.add_slide(blank_layout)
    create_header(s_kuadran, "Matriks 4 Kuadran Portofolio Strategis Program Studi S1 USK (2026)")
    img10 = os.path.join(chart_dir, "10_matriks_4_kuadran_prodi_usk.png")
    if os.path.exists(img10):
        s_kuadran.shapes.add_picture(img10, Inches(0.8), Inches(1.5), width=Inches(8.3))
    
    tb_k = s_kuadran.shapes.add_textbox(Inches(9.3), Inches(1.5), Inches(3.4), Inches(5.4))
    tf_k = tb_k.text_frame
    tf_k.word_wrap = True
    p = tf_k.paragraphs[0]
    p.text = "PANDUAN KEPUTUSAN PIMPINAN:"
    p.font.bold = True
    p.font.size = Pt(13.5)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(8)
    add_bullet(tf_k, "Kuadran I (Prima):", "Peminat tinggi & kursi penuh (Farmasi, Informatika, PGSD, Kedokteran, Hukum). Boleh ekspansi kuota & kelas internasional.", font_size=11, space_after=6)
    add_bullet(tf_k, "Kuadran II (Stabil):", "Peminat sedang tapi kursi terisi penuh (Arsitektur, T. Mesin, Pend. Bhs Indo). Jaga kuota tetap stabil.", font_size=11, space_after=6)
    add_bullet(tf_k, "Kuadran III (Kritis):", "Peminat sepi & kursi banyak bolong (Budidaya Perairan, Fisika, PSP, Proteksi Tanaman, THP). WAJIB PANGKAS KUOTA 20-40%!", font_size=11, space_after=6)
    add_bullet(tf_k, "Kuadran IV (Dilema):", "Peminat tinggi tapi keterisian bocor (Manajemen FEB, Eko. Pembangunan). Perbaiki konversi & skema cicilan IPI.", font_size=11, space_after=6)

    # -------------------------------------------------------------
    # SLIDE 6: TOP TREN PENINGKATAN (GRAFIK 2)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    create_header(s5, "Bintang Pertumbuhan: Program Studi S1 dengan Laju Peningkatan Tertinggi")

    img2 = os.path.join(chart_dir, "02_top_tren_peningkatan_pendaftar_dan_peminat.png")
    if os.path.exists(img2):
        s5.shapes.add_picture(img2, Inches(0.8), Inches(1.6), width=Inches(8.0))
    
    tb = s5.shapes.add_textbox(Inches(9.1), Inches(1.6), Inches(3.6), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ANALISIS TREN TUMBUH:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(46, 125, 50)
    p.space_after = Pt(10)
    add_bullet(tf, "1. Ilmu Keperawatan:", "Slope +38,4 mhs/thn (CAGR +23,5%). Melonjak dari 127 ke 295 mhs. Permintaan lulusan ners internasional ke Jepang & Jerman sangat kuat.")
    add_bullet(tf, "2. PGSD (FKIP):", "Slope +27,4 mhs/thn (CAGR +19,9%). Rekor 4 tahun berturut-turut selalu naik (105 ke 217 mhs) berkat kepastian formasi PPPK Guru.")
    add_bullet(tf, "3. Teknik Sipil & Komputer:", "Masing-masing tumbuh konsisten (Sipil +23 mhs/thn, Komputer +14,8 mhs/thn) sejalan dengan digitalisasi dan infrastruktur.")
    add_bullet(tf, "4. Kehutanan:", "Tumbuh konsisten 4 tahun beruntun (+14,9% CAGR) berkat tren isu lingkungan dan sertifikasi karbon.")

    # -------------------------------------------------------------
    # SLIDE 6: TOP TREN PENURUNAN (GRAFIK 3)
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    create_header(s6, "Sinyal Waspada: Program Studi yang Mengalami Penurunan Berkelanjutan")
    img3 = os.path.join(chart_dir, "03_top_tren_penurunan_pendaftar_dan_peminat.png")
    if os.path.exists(img3):
        s6.shapes.add_picture(img3, Inches(0.8), Inches(1.6), width=Inches(8.0))
    
    tb = s6.shapes.add_textbox(Inches(9.1), Inches(1.6), Inches(3.6), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ANALISIS KONTRAKSI PRODI:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(10)
    add_bullet(tf, "1. Tiga Serangkai FEB Tertekan:", "Manajemen (Slope -16,4 mhs/thn), Ekonomi Islam (-10,3 mhs/thn), dan Ekonomi Pembangunan (-7,7 mhs/thn) mengalami penurunan 3 tahun berturut-turut.")
    add_bullet(tf, "2. Faktor Kanibalisasi Internal:", "Lahirnya prodi baru S1 Bisnis Digital di FEB menyedot langsung calon mahasiswa Manajemen.")
    add_bullet(tf, "3. Sains Murni (Matematika & Fisika):", "Matematika turun -2,9 mhs/thn (CAGR -4,4%). Generasi muda menghindari sains murni karena persepsi minimnya formasi kerja kantoran.")
    add_bullet(tf, "4. Pemanfaatan Sumberdaya Perikanan:", "Turun konsisten dari 80 ke 67 mhs.")

    # -------------------------------------------------------------
    # SLIDE 7: FENOMENA OVER-EKSPANSI KUOTA (GRAFIK 5)
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    create_header(s7, "Fenomena Over-Ekspansi Kuota: Tambah Kuota Tidak Menambah Mahasiswa")
    img5 = os.path.join(chart_dir, "05_over_ekspansi_kuota_vs_daftar_ulang_riil.png")
    if os.path.exists(img5):
        s7.shapes.add_picture(img5, Inches(0.8), Inches(1.6), width=Inches(8.0))
    
    tb = s7.shapes.add_textbox(Inches(9.1), Inches(1.6), Inches(3.6), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "BUKTI OVER-EKSPANSI:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(10)
    add_bullet(tf, "Budidaya Perairan:", "Kuota dibuka 160 kursi, tapi mahasiswa daftar ulang selalu stagnan di angka 90 mhs -> 70 Kursi Kosong!")
    add_bullet(tf, "Teknologi Hasil Pertanian:", "Kuota dinaikkan 2x lipat (80 ke 160 kursi), pendaftar ulang hanya 102 mhs -> 58 Kursi Kosong!")
    add_bullet(tf, "Pendidikan Ekonomi:", "Kuota dinaikkan ke 160 kursi, pendaftar ulang hanya 103 mhs -> 57 Kursi Kosong!")
    add_bullet(tf, "Teknik Kimia:", "Kuota dinaikkan ke 180 kursi, pendaftar ulang hanya 120 mhs -> 60 Kursi Kosong!")
    add_bullet(tf, "Rekomendasi Kebijakan:", "Pangkas kuota ke angka riil 90–120 untuk menghapus ~245 kursi kosong semu.")

    # -------------------------------------------------------------
    # SLIDE 8: DINAMIKA JALUR MASUK & KEBOCORAN (GRAFIK 6)
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    create_header(s8, "Dinamika Jalur Masuk & Titik Kebocoran Calon Mahasiswa (2026)")
    img6 = os.path.join(chart_dir, "06_dinamika_jalur_masuk_dan_kebocoran_2026.png")
    if os.path.exists(img6):
        s8.shapes.add_picture(img6, Inches(0.8), Inches(1.6), width=Inches(8.2))
    
    tb = s8.shapes.add_textbox(Inches(9.2), Inches(1.6), Inches(3.5), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "KONTRIBUSI & KEBOCORAN:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(10)
    add_bullet(tf, "SNBT (38,9%):", "Kontributor terbesar (3.260 mhs). Bocor 626 calon karena diterima di kedinasan/PTN Jawa.")
    add_bullet(tf, "SNBP (32,9%):", "Jalur paling loyal (2.754 mhs). Yield 94,8%, hanya 151 orang gugur berkat sanksi blacklist nasional.")
    add_bullet(tf, "SMMPTN Mandiri (21,6%):", "1.814 mhs masuk, tapi 604 orang (25%) gugur akibat shock tagihan IPI belasan juta dalam tempo 5 hari.")
    add_bullet(tf, "TALENTA (Bocor Terparah):", "930 calon gugur (Yield hanya 24,9%). 3 dari 4 orang menjadikan USK cadangan gratis.")

    # -------------------------------------------------------------
    # SLIDE 9: KINERJA KETERISIAN PER FAKULTAS (GRAFIK 7)
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    create_header(s9, "Evaluasi Kinerja 12 Fakultas: Keterisian Kuota 2022 vs 2026")
    img7 = os.path.join(chart_dir, "07_analisa_peminatan_dan_keterisian_fakultas.png")
    if os.path.exists(img7):
        s9.shapes.add_picture(img7, Inches(0.8), Inches(1.6), width=Inches(7.8))
    
    tb = s9.shapes.add_textbox(Inches(8.9), Inches(1.6), Inches(3.8), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ZONASI KINERJA FAKULTAS:"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(10)
    add_bullet(tf, "Zona Prima (98-100%):", "FKG (100%), Fakultas Hukum (99,6%), dan Fakultas Kedokteran (98,8%). Kursi selalu terisi penuh sempurna.")
    add_bullet(tf, "Zona Sehat (>80%):", "FISIP (89,6%), FKH (88,5%), Keperawatan (81,9%), FMIPA (81,9%), FKIP (81,5%), FT (81,4%).")
    add_bullet(tf, "Zona Rentan & Kritis (<80%):", "FEB (77,2%), FPK (62,4%), dan Fakultas Pertanian (62,1%). Perlu intervensi rasionalisasi kuota segera.")

    # -------------------------------------------------------------
    # SLIDE 10: KLASTER KHUSUS D3 & PSDKU (GRAFIK 8 & 9)
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    create_header(s10, "Evaluasi Klaster Khusus: Krisis Vokasi D3 & Outlier PSDKU Gayo Lues")
    img8 = os.path.join(chart_dir, "08_evaluasi_multi_tahun_d3_vokasi.png")
    img9 = os.path.join(chart_dir, "09_subanalisis_psdku_gayo_lues.png")
    if os.path.exists(img8):
        s10.shapes.add_picture(img8, Inches(0.8), Inches(1.6), width=Inches(5.7))
    if os.path.exists(img9):
        s10.shapes.add_picture(img9, Inches(6.8), Inches(1.6), width=Inches(5.7))

    # -------------------------------------------------------------
    # SLIDE 11: EKSPLORASI KAUSALITAS LAPANGAN
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    create_header(s11, "Eksplorasi Kausalitas Lapangan: Membedah Fakta Data vs Hipotesis")
    
    # 2 Big Cards: FAKTA DATA vs HIPOTESIS
    c_fakta = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.3))
    set_shape_flat_color(c_fakta, RGBColor(240, 244, 248))
    c_fakta.line.color.rgb = RGBColor(27, 59, 111)
    tf_f = c_fakta.text_frame
    tf_f.word_wrap = True
    tf_f.margin_left = tf_f.margin_right = tf_f.margin_top = Inches(0.25)
    p = tf_f.paragraphs[0]
    p.text = "[FAKTA DATA EMPIRIS RESMI]"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(12)
    add_bullet(tf_f, "Fakta 1:", "PGSD naik +106% (105 ke 217 mhs) dan Keperawatan naik +132% (127 ke 295 mhs).")
    add_bullet(tf_f, "Fakta 2:", "Manajemen merosot 3 tahun beruntun (-25%), bertepatan dengan dibukanya prodi Bisnis Digital di tahun 2024.")
    add_bullet(tf_f, "Fakta 3:", "Budidaya Perairan memiliki peminat 145 orang untuk 160 kursi (keketatan < 1:1) dan 70 kursi kosong di 2026.")
    add_bullet(tf_f, "Fakta 4:", "25% calon lulus Jalur Mandiri (604 orang) dan 75% Jalur Talenta (930 orang) mengundurkan diri.")
    add_bullet(tf_f, "Fakta 5:", "4 prodi PSDKU Gayo Lues hanya terisi 41 mahasiswa dari 220 kursi (81,4% kosong).")

    c_hipo = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.7), Inches(5.3))
    set_shape_flat_color(c_hipo, RGBColor(254, 249, 231))
    c_hipo.line.color.rgb = RGBColor(230, 81, 0)
    tf_h = c_hipo.text_frame
    tf_h.word_wrap = True
    tf_h.margin_left = tf_h.margin_right = tf_h.margin_top = Inches(0.25)
    p = tf_h.paragraphs[0]
    p.text = "[TEMUAN DESK RESEARCH & PENYEBAB]"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(183, 121, 31)
    p.space_after = Pt(12)
    add_bullet(tf_h, "Faktor PPPK & Ners Luar Negeri:", "Formasi ASN PPPK Guru SD di Aceh dan program G-to-G perawat ke Jepang & Jerman memberikan jaminan karir instan.")
    add_bullet(tf_h, "Faktor Kanibalisasi Internal:", "Bisnis Digital menyedot segmen pendaftar Manajemen yang menginginkan gelar kekinian di era ekonomi digital.")
    add_bullet(tf_h, "Faktor Persepsi Karir Agromaritim:", "Gen-Z menghindari prodi perikanan/sains murni karena dipersepsikan minim lowongan kerja korporasi di daerah.")
    add_bullet(tf_h, "Beban Finansial IPI Mandiri:", "Kewajiban melunasi IPI belasan hingga puluhan juta dalam 5 hari (SK Rektor 1162/2026) memicu liquidity shock keluarga.")
    add_bullet(tf_h, "Faktor Isolasi Geografis:", "Akses darat 10-12 jam ke Gayo Lues membatasi pendaftar hanya pada lulusan SMA setempat.")

    # -------------------------------------------------------------
    # SLIDE 12: INTEGRASI JALUR INTERNASIONAL & FLYER PMB
    # -------------------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    create_header(s12, "Integrasi Jalur Internasional & Rekomendasi Media Promosi Flyer PMB")
    
    tb = s12.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PEMUTAKHIRAN INFORMASI SESUAI BUTIR 17 PANDUAN MAGANG:"
    p.font.bold = True
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(27, 59, 111)
    p.space_after = Pt(15)

    add_bullet(tf, "1. Peluang Mahasiswa Internasional (International Student Admissions):", "USK memiliki daya tarik kuat bagi mahasiswa asal Malaysia, Thailand Selatan, dan Timur Tengah pada prodi Kedokteran, Kedokteran Hewan, dan Kebencanaan (TDMRC). Kuota internasional dapat dimanfaatkan untuk mengisi kelas unggulan dengan standar UKT internasional.")
    add_bullet(tf, "2. Pembaruan Materi Promosi (Flyer PMB 2027):", "Materi promosi tidak boleh lagi bersifat umum satu brosur untuk semua. Perlu dibuat Flyer Tersegmentasi: (a) Flyer Vokasi Unggulan menonjolkan sertifikasi kompetensi, (b) Flyer Beasiswa Ikatan Dinas Pemda untuk PSDKU Gayo Lues, dan (c) Flyer Program Ners Global untuk Keperawatan.")
    add_bullet(tf, "3. Transparansi Simulasi Biaya UKT & IPI:", "Flyer resmi PMB wajib menyertakan simulasi tabel IPI dan fasilitas cicilan secara transparan agar orang tua calon mahasiswa tidak mengalami shock saat pengumuman kelulusan mandiri.")

    # -------------------------------------------------------------
    # SLIDE 13: REKOMENDASI KEBIJAKAN RASIONALISASI KUOTA 2027
    # -------------------------------------------------------------
    s13 = prs.slides.add_slide(blank_layout)
    create_header(s13, "Rekomendasi Preskriptif PMB 2027: Rasionalisasi Kuota Menuju Efisiensi PTN-BH")
    
    # 4 Action Cards
    w4 = Inches(2.7)
    y4 = Inches(1.6)
    h4 = Inches(5.2)

    # Box 1
    b1 = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y4, w4, h4)
    set_shape_flat_color(b1, RGBColor(255, 255, 255))
    b1.line.color.rgb = RGBColor(192, 57, 43)
    tf1 = b1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "AKSI 1:\nPANGKAS KUOTA"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(192, 57, 43)
    p.space_after = Pt(10)
    add_bullet(tf1, "7 Prodi Kritis:", "Pangkas kuota 20-40% pada: Budidaya Perairan (160 -> 90), THP (160 -> 100), Pend. Ekonomi (160 -> 100), Teknik Kimia (180 -> 120), Fisika (80 -> 50), dan PSP (120 -> 70).")
    add_bullet(tf1, "Dampak Langsung:", "Menghapus ~400 kursi kosong semu dan mendongkrak rasio pemenuhan USK di atas 85%.")

    # Box 2
    b2 = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), y4, w4, h4)
    set_shape_flat_color(b2, RGBColor(255, 255, 255))
    b2.line.color.rgb = RGBColor(230, 81, 0)
    tf2 = b2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "AKSI 2:\nREFORMASI JALUR"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(230, 81, 0)
    p.space_after = Pt(10)
    add_bullet(tf2, "Kunci Jalur Talenta:", "Terapkan uang komitmen registrasi Rp1.000.000 (memotong UKT jika masuk, hangus jika kabur) agar tidak jadi cadangan gratis.")
    add_bullet(tf2, "Cicilan IPI Mandiri:", "Buka skema cicilan 3 tahap sepanjang semester 1 & 2 untuk menyelamatkan 600+ calon mahasiswa yang terkendala biaya mendadak.")

    # Box 3
    b3 = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), y4, w4, h4)
    set_shape_flat_color(b3, RGBColor(255, 255, 255))
    b3.line.color.rgb = RGBColor(41, 128, 185)
    tf3 = b3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "AKSI 3:\nUPGRADE VOKASI"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(41, 128, 185)
    p.space_after = Pt(10)
    add_bullet(tf3, "Migrasi ke D4:", "Konversi D3 favorit ke Sarjana Terapan (D4): D3 Manajemen Informatika -> D4 Sains Data Terapan; D3 Sipil -> D4 Manajemen Konstruksi.")
    add_bullet(tf3, "Daya Tarik ASN:", "Ijazah D4 setara S1 berhak atas golongan III/a ASN, menyelesaikan krisis peminat D3.")

    # Box 4
    b4 = s13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), y4, w4, h4)
    set_shape_flat_color(b4, RGBColor(255, 255, 255))
    b4.line.color.rgb = RGBColor(39, 174, 96)
    tf4 = b4.text_frame
    tf4.word_wrap = True
    p = tf4.paragraphs[0]
    p.text = "AKSI 4:\nSOLUSI PSDKU"
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(39, 174, 96)
    p.space_after = Pt(10)
    add_bullet(tf4, "MoU Beasiswa Pemkab:", "Ikat 220 kuota Gayo Lues dengan beasiswa penuh APBD Pemkab Gayo Lues & Aceh Tenggara.")
    add_bullet(tf4, "Rampingkan Operasional:", "Jika beasiswa tidak terwujud, konsolidasi 4 prodi menjadi 2 prodi unggulan lokal: Agribisnis Kopi Gayo dan Keguruan.")

    # -------------------------------------------------------------
    # SLIDE 14: PENUTUP
    # -------------------------------------------------------------
    s14 = prs.slides.add_slide(blank_layout)
    bg14 = s14.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_shape_flat_color(bg14, RGBColor(27, 59, 111))

    tb = s14.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "TERIMA KASIH"
    p.font.bold = True
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

    p2 = tf.add_paragraph()
    p2.text = "\"Keberhasilan PMB PTN-BH diukur dari berapa banyak mahasiswa yang secara nyata duduk di ruang kuliah, membayar UKT, dan menyelesaikan studinya secara tepat waktu, bukan dari angka kelulusan semu di atas kertas.\""
    p2.font.size = Pt(16)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(230, 230, 230)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(25)

    p3 = tf.add_paragraph()
    p3.text = "Universitas Syiah Kuala | Kampus Jantong Hate Rakyat Aceh"
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(230, 81, 0)
    p3.alignment = PP_ALIGN.CENTER

    prs.save(out_pptx)
    print(f"Presentation PPTX successfully saved to: {out_pptx}")

if __name__ == "__main__":
    main()
